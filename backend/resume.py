from fastapi import APIRouter, UploadFile, File, Depends, HTTPException
from secure import verify_token
import requests
import fitz
import os
from dotenv import load_dotenv
import openpyxl
from io import BytesIO
from pptx import Presentation
import docx
from models import history, users
from databases import Database


load_dotenv()

HF_API_TOKEN = os.getenv("HF_API_TOKEN")
HF_MODEL = "facebook/bart-large-cnn"
HF_API_URL = f"https://api-interface.huggingface.co/models/{HF_MODEL}"

router = APIRouter()

MAX_CHUNKS_PER_FILE = 5
CHUNCK_SIZE = 2000

def extract_text_from_pdf(file_bytes:bytes)->str:
    text = ""
    with fitz.open(stream=file_bytes, filetype="pdf") as doc:
        for page in doc:
            text += page.get_text()
    return text

def extract_text_from_excel(file_bytes:bytes)->str:
    wb = openpyxl.load_workbook(filename=BytesIO(file_bytes), data_only=True)
    text = ""
    for sheet in wb.worksheets:
        for row in sheet.iter_rows(values_only=True):
            for cell in row:
                if cell is not None:
                    text += str(cell) + " "
    return text

def extract_text_from_pptx(file_bytes:bytes)->str:
    prs = Presentation(BytesIO(file_bytes))
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + " "
    return text

def extract_text_from_docx(file_bytes:bytes)->str:
    doc = docx.Document(BytesIO(file_bytes))
    text = ""
    for para in doc.paragraphs:
        text += para.text + " "
    return text
    
def call_hf_api(text:str)->str:
    headers = {"Authorization": f"Bearer {HF_API_TOKEN}"}
    payload = {"inputs": text, "options": {"wait_for_model": True}}
    response = requests.post(HF_API_URL, headers=headers, json=payload)
    if response.status_code != 200:
        raise HTTPException(status_code=500, detail=f"Erreur lors de l'appel à l'API Hugging Face: {response.text}")
    result = response.json()
    if isinstance(result, list) and "summary_text" in result[0]:
        return result[0]["summary_text"]
    raise HTTPException(status_code=500, detail="Erreur lors de l'appel à l'API Hugging Face: Format de réponse inattendue")

@router.post("/upload")
async def upload_file(file: UploadFile = File(...), user=Depends(verify_token)):
    content = await file.read()
    text = ""
    if file.filename.endswith(".txt"):
        text = content.decode("utf-8")
    elif file.filename.endswith(".pdf"):
        text = extract_text_from_pdf(content)
    elif file.filename.endswith(".xlsx") or file.filename.endswith(".xls"):
        text = extract_text_from_excel(content)
    elif file.filename.endswith(".pptx"):
        text = extract_text_from_pptx(content)
    elif file.filename.endswith(".docx"):
        text = extract_text_from_docx(content)
    else:
        raise HTTPException(status_code=400, detail="Format de fichier non supporté. Seuls les fichiers .txt, .pdf, .xlsx, .pptx et .docx sont acceptés.")
    
    if len(text.strip()) == 0:
        raise HTTPException(status_code=400, detail="Le fichier est vide ou le texte n'a pas pu être extrait.")
    
    chunks = [text[i:i+CHUNCK_SIZE] for i in range(0, len(text), CHUNCK_SIZE)]
    if len(chunks) > MAX_CHUNKS_PER_FILE:
        chunks = chunks[:MAX_CHUNKS_PER_FILE]
        
    summaries = []
    for chunk in chunks:
        summary = call_hf_api(chunk)
        summaries.append(summary)
    final_summary = " ".join(summaries)
    
    return {
        "user": user,
        "filename": file.filename,
        "summary": final_summary,
        "chunks_used": len(chunks)
    }
    
    user_query = users.select().where(users.c.email == user_email)
    user_db = await database.fetch_one(user_query)
    user_id = user_db.id if user_db else None
    
    insert_query = history.insert().values(
        user_id=user_id,
        file_name=file.filename,
        file_type=file.filename.split(".")[-1],
        format=format,
        summary=final_summary
    )
    
    await database.execute(insert_query)